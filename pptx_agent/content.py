from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from .models import ContentPlan, SlideContentUpdate


def _trim_line(text: str, max_len: int = 220) -> str:
    clean = re.sub(r"\s+", " ", text).strip()
    if len(clean) <= max_len:
        return clean
    return clean[: max_len - 1].rstrip() + "..."


def _find_title_shape(slide):
    for shape in slide.shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        try:
            ph_type = shape.placeholder_format.type
        except ValueError:
            continue
        if ph_type in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}:
            return shape
    for shape in slide.shapes:
        if shape.has_text_frame:
            return shape
    return None


def _find_subtitle_shape(slide):
    for shape in slide.shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        try:
            ph_type = shape.placeholder_format.type
        except ValueError:
            continue
        if ph_type == PP_PLACEHOLDER.SUBTITLE:
            return shape
    return None


def _body_text_shapes(slide):
    body = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if getattr(shape, "is_placeholder", False):
            try:
                ph_type = shape.placeholder_format.type
            except ValueError:
                ph_type = None
            if ph_type in {
                PP_PLACEHOLDER.TITLE,
                PP_PLACEHOLDER.CENTER_TITLE,
                PP_PLACEHOLDER.SUBTITLE,
            }:
                continue
        body.append(shape)
    return body


def _write_line(paragraph, line: str, bold_all: bool = False) -> None:
    paragraph.clear()
    cleaned = _trim_line(line)
    if not cleaned:
        return

    if bold_all:
        run = paragraph.add_run()
        run.text = cleaned
        run.font.bold = True
        return

    label_match = re.match(r"^([^:]{1,80}:)\s*(.+)$", cleaned)
    if label_match:
        label_run = paragraph.add_run()
        label_run.text = label_match.group(1)
        label_run.font.bold = True

        value = label_match.group(2).strip()
        if value:
            value_run = paragraph.add_run()
            value_run.text = f" {value}"
        return

    run = paragraph.add_run()
    run.text = cleaned


def _write_text_frame(shape, lines: list[str], bullets: bool) -> None:
    if not shape.has_text_frame:
        return
    frame = shape.text_frame
    frame.clear()
    if not lines:
        return

    for idx, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.level = 0
        _write_line(paragraph, line, bold_all=False)

        if not bullets:
            continue


def _bold_all_runs(shape) -> None:
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        if not paragraph.runs and paragraph.text.strip():
            run = paragraph.add_run()
            run.text = paragraph.text
        for run in paragraph.runs:
            run.font.bold = True


def _remove_empty_placeholders(slide) -> int:
    to_remove = []
    for shape in slide.shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        if not shape.has_text_frame:
            continue
        text = shape.text.strip()
        if text:
            continue
        try:
            ph_type = shape.placeholder_format.type
        except ValueError:
            continue
        if ph_type in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}:
            continue
        to_remove.append(shape)

    for shape in to_remove:
        element = shape.element
        parent = element.getparent()
        if parent is not None:
            parent.remove(element)
    return len(to_remove)


def apply_slide_update(slide, update: SlideContentUpdate) -> dict[str, int]:
    counters = {"removed_placeholders": 0}

    title_shape = _find_title_shape(slide)
    if update.title and title_shape is not None:
        title_shape.text = _trim_line(update.title, max_len=180)
        _bold_all_runs(title_shape)

    subtitle_shape = _find_subtitle_shape(slide)
    if update.subtitle and subtitle_shape is not None:
        subtitle_shape.text = _trim_line(update.subtitle, max_len=220)

    body_shapes = _body_text_shapes(slide)
    cursor = 0

    if update.bullets and body_shapes:
        _write_text_frame(body_shapes[cursor], update.bullets[:10], bullets=True)
        cursor += 1

    if update.body_paragraphs:
        target = (
            body_shapes[cursor]
            if cursor < len(body_shapes)
            else (body_shapes[-1] if body_shapes else None)
        )
        if target is not None:
            _write_text_frame(target, update.body_paragraphs[:8], bullets=False)

    if update.notes:
        try:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_frame.text = _trim_line(update.notes, max_len=400)
        except Exception:
            pass

    if update.remove_empty_placeholders:
        counters["removed_placeholders"] = _remove_empty_placeholders(slide)

    return counters


def apply_content_plan(
    input_pptx: Path, content_plan: ContentPlan, output_pptx: Path
) -> dict[str, int]:
    prs = Presentation(str(input_pptx))
    totals = {"slides_updated": 0, "removed_placeholders": 0}

    for update in content_plan.slides:
        if update.slide_index < 1 or update.slide_index > len(prs.slides):
            continue
        slide = prs.slides[update.slide_index - 1]
        counters = apply_slide_update(slide, update)
        totals["slides_updated"] += 1
        totals["removed_placeholders"] += counters["removed_placeholders"]

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_pptx))
    return totals
