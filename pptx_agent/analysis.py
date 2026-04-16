from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation

from .models import DeckAnalysis, SlideSummary
from .qa import render_slides_to_images
from .utils import run_command


def extract_text_markitdown(pptx_path: Path) -> str:
    result = run_command([sys.executable, "-m", "markitdown", str(pptx_path)])
    if result.returncode != 0:
        return ""
    return result.stdout.strip()


def extract_text_python_pptx(pptx_path: Path) -> str:
    prs = Presentation(str(pptx_path))
    lines: list[str] = []
    for index, slide in enumerate(prs.slides, start=1):
        lines.append(f"# Slide {index}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    lines.append(text)
        lines.append("")
    return "\n".join(lines).strip()


def detect_language(text: str, default: str = "fr") -> str:
    if not text:
        return default
    lowered = text.lower()
    french_markers = [" le ", " la ", " les ", " de ", " pour ", " avec "]
    english_markers = [" the ", " and ", " for ", " with ", " is ", " are "]
    fr_score = sum(marker in lowered for marker in french_markers)
    en_score = sum(marker in lowered for marker in english_markers)
    if en_score > fr_score:
        return "en"
    return "fr"


def _shape_text_preview(shape_text: str, max_len: int = 180) -> str:
    clean = re.sub(r"\s+", " ", shape_text).strip()
    if len(clean) <= max_len:
        return clean
    return clean[: max_len - 1].rstrip() + "..."


def _collect_slide_summary(pptx_path: Path) -> list[SlideSummary]:
    prs = Presentation(str(pptx_path))
    slides: list[SlideSummary] = []
    for index, slide in enumerate(prs.slides, start=1):
        placeholders: list[str] = []
        text_blocks: list[str] = []

        for shape in slide.shapes:
            if getattr(shape, "is_placeholder", False):
                placeholders.append(shape.name)
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    text_blocks.append(text)

        layout_name = None
        if slide.slide_layout is not None:
            layout_name = slide.slide_layout.name

        slides.append(
            SlideSummary(
                slide_index=index,
                layout_name=layout_name,
                placeholder_names=placeholders,
                text_preview=_shape_text_preview(" | ".join(text_blocks)),
            )
        )

    return slides


def analyze_template(
    pptx_path: Path, workdir: Path, default_language: str = "fr"
) -> DeckAnalysis:
    workdir.mkdir(parents=True, exist_ok=True)
    thumbnails_dir = workdir / "thumbnails"
    thumbnails_dir.mkdir(parents=True, exist_ok=True)

    extracted = extract_text_markitdown(pptx_path)
    if not extracted:
        extracted = extract_text_python_pptx(pptx_path)

    slides = _collect_slide_summary(pptx_path)

    thumbnail_paths: list[str] = []
    render_notes: list[str] = []
    try:
        images = render_slides_to_images(
            pptx_path, thumbnails_dir, prefix="analysis-slide"
        )
        thumbnail_paths = [str(path) for path in images]
    except RuntimeError as exc:
        render_notes.append(str(exc))

    language = detect_language(extracted, default=default_language)

    analysis = DeckAnalysis(
        source_path=str(pptx_path),
        slide_count=len(slides),
        detected_language=language,
        extracted_text=extracted,
        thumbnail_paths=thumbnail_paths,
        slides=slides,
    )

    text_path = workdir / "analysis.txt"
    payload = extracted
    if render_notes:
        payload += "\n\n[notes]\n" + "\n".join(render_notes)
    text_path.write_text(payload, encoding="utf-8")

    return analysis
