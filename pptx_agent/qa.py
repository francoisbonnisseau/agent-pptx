from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation

from .models import QAReport, QAIssue, QAIssueCategory, QAIssueSeverity
from .utils import command_exists, extract_json_object, natural_sort_key, run_command

PLACEHOLDER_REGEX = re.compile(
    r"\b(x{3,}|lorem|ipsum|placeholder|todo|tbd|this\s+.*(?:slide|page).*layout)\b",
    flags=re.IGNORECASE,
)


def _extract_text_markitdown(pptx_path: Path) -> str:
    result = run_command([sys.executable, "-m", "markitdown", str(pptx_path)])
    if result.returncode != 0:
        return ""
    return result.stdout


def _extract_text_python_pptx(pptx_path: Path) -> str:
    prs = Presentation(str(pptx_path))
    lines: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        lines.append(f"# Slide {idx}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                lines.append(shape.text.strip())
    return "\n".join(lines)


def run_content_qa(pptx_path: Path) -> QAReport:
    text = _extract_text_markitdown(pptx_path)
    if not text:
        text = _extract_text_python_pptx(pptx_path)

    issues: list[QAIssue] = []

    current_slide = 1
    for raw_line in text.splitlines():
        line = raw_line.strip()
        match_slide = re.search(r"slide\s+(\d+)", line, flags=re.IGNORECASE)
        if match_slide:
            current_slide = max(1, int(match_slide.group(1)))

        if PLACEHOLDER_REGEX.search(line):
            issues.append(
                QAIssue(
                    slide_index=current_slide,
                    severity=QAIssueSeverity.high,
                    category=QAIssueCategory.placeholder,
                    description=f"Texte placeholder detecte: {line[:180]}",
                    fix_hint="Remplacer par le contenu final et supprimer les elements inutilises.",
                )
            )

        if "\u2022" in line:
            issues.append(
                QAIssue(
                    slide_index=current_slide,
                    severity=QAIssueSeverity.medium,
                    category=QAIssueCategory.other,
                    description="Bullet Unicode detecte (\\u2022).",
                    fix_hint="Utiliser des paragraphes de liste natifs PowerPoint.",
                )
            )

    prs = Presentation(str(pptx_path))
    for idx, slide in enumerate(prs.slides, start=1):
        has_text = False
        has_title_like = False
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                has_text = True
                if len(shape.text.strip()) <= 120 and "\n" not in shape.text:
                    has_title_like = True

        if not has_text:
            issues.append(
                QAIssue(
                    slide_index=idx,
                    severity=QAIssueSeverity.high,
                    category=QAIssueCategory.other,
                    description="Slide vide detectee.",
                    fix_hint="Ajouter du contenu ou supprimer la slide.",
                )
            )
        elif not has_title_like:
            issues.append(
                QAIssue(
                    slide_index=idx,
                    severity=QAIssueSeverity.low,
                    category=QAIssueCategory.other,
                    description="Aucun titre evident detecte.",
                    fix_hint="Verifier qu'un titre clair est present.",
                )
            )

    return QAReport(mode="content", issues=issues)


def _resolve_soffice_command() -> str | None:
    for candidate in ("soffice", "libreoffice"):
        if command_exists(candidate):
            return candidate
    return None


def _resolve_pdftoppm_command() -> str | None:
    if command_exists("pdftoppm"):
        return "pdftoppm"
    return None


def render_slides_to_images(
    pptx_path: Path, output_dir: Path, prefix: str = "slide"
) -> list[Path]:
    soffice_cmd = _resolve_soffice_command()
    pdftoppm_cmd = _resolve_pdftoppm_command()
    if not soffice_cmd or not pdftoppm_cmd:
        raise RuntimeError("soffice et/ou pdftoppm indisponibles dans PATH")

    output_dir.mkdir(parents=True, exist_ok=True)
    pdf_path = output_dir / f"{pptx_path.stem}.pdf"

    convert = run_command(
        [
            soffice_cmd,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir),
            str(pptx_path),
        ]
    )
    if convert.returncode != 0:
        raise RuntimeError(
            f"Conversion PDF impossible: {convert.stderr.strip() or convert.stdout.strip()}"
        )

    if not pdf_path.exists():
        fallback_pdf = output_dir / f"{pptx_path.name}.pdf"
        if fallback_pdf.exists():
            pdf_path = fallback_pdf
        else:
            raise RuntimeError("PDF intermediaire introuvable apres conversion")

    render = run_command(
        [pdftoppm_cmd, "-png", "-r", "150", str(pdf_path), str(output_dir / prefix)]
    )
    if render.returncode != 0:
        raise RuntimeError(
            f"Rendu PNG impossible: {render.stderr.strip() or render.stdout.strip()}"
        )

    images = sorted(output_dir.glob(f"{prefix}-*.png"), key=natural_sort_key)
    if not images:
        images = sorted(output_dir.glob(f"{prefix}*.png"), key=natural_sort_key)
    if not images:
        raise RuntimeError("Aucune image de slide generee")
    return images


def _parse_slide_index_from_path(path: Path) -> int:
    match = re.search(r"(\d+)(?=\.png$)", path.name)
    if not match:
        return 1
    return max(1, int(match.group(1)))


def run_visual_qa_with_gemini(
    pptx_path: Path,
    output_dir: Path,
    gemini_api_key: str | None,
    model_name: str,
) -> QAReport:
    notes: list[str] = []

    if not gemini_api_key:
        notes.append("GEMINI_API_KEY absent, visual QA ignoree.")
        return QAReport(mode="visual", notes=notes)

    if not command_exists("pdftoppm") or (_resolve_soffice_command() is None):
        notes.append("soffice/pdftoppm indisponibles, visual QA ignoree.")
        return QAReport(mode="visual", notes=notes)

    try:
        import google.generativeai as genai
        from PIL import Image
    except Exception:
        notes.append("google.generativeai ou Pillow indisponible, visual QA ignoree.")
        return QAReport(mode="visual", notes=notes)

    images = render_slides_to_images(pptx_path, output_dir, prefix="qa-slide")
    notes.append(f"{len(images)} slide(s) rendues pour visual QA.")

    genai.configure(api_key=gemini_api_key)
    model = genai.GenerativeModel(model_name)

    issues: list[QAIssue] = []
    for image_path in images:
        slide_idx = _parse_slide_index_from_path(image_path)
        prompt = (
            "Analyse cette slide de presentation et retourne strictement un JSON de la forme "
            '{"issues":[{"severity":"high|medium|low","category":"overflow|overlap|alignment|contrast|spacing|other",'
            '"description":"...","fix_hint":"..."}]}. '
            'Si tout est correct, retourne {"issues":[]}. '
            "Cherche explicitement: overlaps, clipping, debordement texte, contrastes faibles, alignements incoherents,"
            " elements trop proches des bords."
        )

        with Image.open(image_path) as image:
            response = model.generate_content(
                [prompt, image],
                generation_config={"temperature": 0.0},
            )

        raw = getattr(response, "text", "") or ""
        parsed = extract_json_object(raw)
        parsed_issues = parsed.get("issues") if isinstance(parsed, dict) else None

        if not isinstance(parsed_issues, list):
            lowered = raw.lower()
            if "aucun" in lowered or "no issue" in lowered or 'issues":[]' in lowered:
                continue
            issues.append(
                QAIssue(
                    slide_index=slide_idx,
                    severity=QAIssueSeverity.low,
                    category=QAIssueCategory.other,
                    description="Reponse vision non parsee en JSON.",
                    fix_hint="Relancer la QA visuelle pour confirmation.",
                )
            )
            continue

        for item in parsed_issues:
            if not isinstance(item, dict):
                continue
            sev_raw = str(item.get("severity", "medium")).lower().strip()
            cat_raw = str(item.get("category", "other")).lower().strip()

            severity = {
                "high": QAIssueSeverity.high,
                "medium": QAIssueSeverity.medium,
                "low": QAIssueSeverity.low,
            }.get(sev_raw, QAIssueSeverity.medium)

            category = {
                "overflow": QAIssueCategory.overflow,
                "overlap": QAIssueCategory.overlap,
                "alignment": QAIssueCategory.alignment,
                "contrast": QAIssueCategory.contrast,
                "spacing": QAIssueCategory.spacing,
                "placeholder": QAIssueCategory.placeholder,
                "typo": QAIssueCategory.typo,
                "other": QAIssueCategory.other,
            }.get(cat_raw, QAIssueCategory.other)

            issues.append(
                QAIssue(
                    slide_index=slide_idx,
                    severity=severity,
                    category=category,
                    description=str(item.get("description", "Issue visuel detecte"))[
                        :300
                    ],
                    fix_hint=str(item.get("fix_hint", "Ajuster la mise en page"))[:300],
                )
            )

    return QAReport(mode="visual", issues=issues, notes=notes)


def merge_issue_counts(reports: list[QAReport]) -> int:
    return sum(len(report.issues) for report in reports)
